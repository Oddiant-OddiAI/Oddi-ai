import base64
import os
import subprocess
import tempfile
import shutil
from app.fast_responses import fast_response, is_job_context_question
from app.memory_handler import (
    is_memory_message,
    remember,
    recall_memory
)
from app.chatbot import get_response
from app.commands import handle_command
from app.prompts import SYSTEM_PROMPT
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document

import cv2

from app.config import client
from app.database import (
    get_vector_store_id,
    save_vector_store_id,
    save_memory
)
from app import state


RESUME_ANALYSIS_TRIGGERS = {
    "analyze my resume",
    "analyse my resume",
    "analyze my cv",
    "analyse my cv",
    "resume analysis",
    "analyze resume",
    "analyse resume",
    "review my resume",
    "review my cv",
}


def is_resume_analysis_request(message):
    if not message:
        return False

    text = message.lower().strip()

    return text in RESUME_ANALYSIS_TRIGGERS
def is_resume_analysis_request(message):
    if not message:
        return False

    text = message.lower().strip()

    return text in RESUME_ANALYSIS_TRIGGERS


def extract_docx_text(uploaded_file):
    document = Document(uploaded_file)

    text = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            text.append(paragraph.text.strip())

    # Also read tables because resumes often contain
    # skills, education, experience, etc. inside tables.
    for table in document.tables:

        for row in table.rows:

            cells = []

            for cell in row.cells:
                value = cell.text.strip()

                if value:
                    cells.append(value)

            if cells:
                text.append(" | ".join(cells))

    return "\n".join(text)
def transcribe_audio(audio_file):
    """
    Transcribe an uploaded audio file using OpenAI's audio transcription API.
    """
    audio_file.seek(0)

    transcription = client.audio.transcriptions.create(
        model="gpt-4o-mini-transcribe",
        file=audio_file
    )

    return transcription.text

RESUME_ANALYSIS_PROMPT = """
You are Oddi AI's Resume Intelligence system.

The user has uploaded a resume and explicitly asked you to analyze it.

Analyze ONLY the resume content provided in the attached document.
Do not invent education, experience, projects, skills, achievements, certifications,
or other information that is not present.

Give the user a practical and honest resume review.

Your response MUST use this structure:

# 📄 Resume Analysis

## Overall Score
Give a score out of 100.

## ATS Score
Give an estimated ATS-readiness score out of 100.
Explain briefly what affects the score.

## 👤 Profile
Summarize the candidate's profile based only on the resume.

## 🎓 Education
Summarize the education shown in the resume.

## 💻 Technical Skills
List the technical skills actually present.

## 🚀 Projects
List and briefly evaluate the projects shown.

## 💼 Experience
Summarize internships, jobs, training, or other experience shown.

If no experience is present, clearly say so.

## ✅ Strengths
Give the strongest aspects of the resume.

## ⚠️ Weaknesses
Give specific weaknesses or areas that reduce its effectiveness.

## ❌ Missing / Recommended Sections
Mention important sections that appear to be missing.
Do not claim a section is missing if it is actually present.

## 🎯 Suitable Roles
Suggest suitable job roles based only on the candidate's demonstrated
education, skills, projects, and experience.

## 📈 Recommended Improvements
Give specific, actionable improvements.
Prioritize the most important changes first.

## 📝 Final Verdict
Give a concise overall assessment.

IMPORTANT RULES:

- Do not invent information.
- Do not change the candidate's facts.
- Do not judge the candidate's personality or appearance.
- Do not make unrealistic promises about getting a job.
- Scores are estimates, not official ATS scores.
- Be honest but constructive.
- Prefer specific feedback over generic advice.
- Keep the report readable and well formatted.
"""


MAX_KNOWLEDGE_FILES = 20


def get_or_create_vector_store(user_id):

    if not user_id:
        return None

    vector_store_id = get_vector_store_id(user_id)

    if vector_store_id:
        return vector_store_id

    vector_store = client.vector_stores.create(
        name=f"Oddi AI User {user_id} Knowledge"
    )

    save_vector_store_id(
        user_id,
        vector_store.id
    )

    print(
        f"Created knowledge base for user {user_id}: "
        f"{vector_store.id}"
    )

    return vector_store.id


def get_knowledge_file_count(vector_store_id):

    files = client.vector_stores.files.list(
        vector_store_id=vector_store_id
    )

    return len(files.data)


def add_file_to_knowledge_base(
    uploaded_file,
    vector_store_id
):

    if not vector_store_id:
        return None

    uploaded_file.stream.seek(0)

    print(
        f"Uploading to knowledge base: "
        f"{uploaded_file.filename}"
    )

    # Upload the original file to OpenAI
    openai_file = client.files.create(
        file=(
            uploaded_file.filename,
            uploaded_file.stream,
            uploaded_file.mimetype
        ),
        purpose="assistants"
    )

    print(
        f"OpenAI file created: {openai_file.id}"
    )

    # Attach it to the user's vector store
    vector_file = client.vector_stores.files.create_and_poll(
        vector_store_id=vector_store_id,
        file_id=openai_file.id
    )

    if vector_file.status != "completed":

        print(
            f"Knowledge indexing failed: "
            f"{vector_file.status}"
        )

        return None

    print(
        f"Indexed successfully: "
        f"{uploaded_file.filename}"
    )

    return openai_file.id
def process_message(
    user_message,
    uploaded_files=None,
    conversation_history=None,
    user_id=None
    ):
        # ==========================================
    # PERSISTENT USER KNOWLEDGE
    # ==========================================

    vector_store_id = None

    if user_id is not None:

        vector_store_id = get_or_create_vector_store(user_id)

        if uploaded_files:

            current_count = get_knowledge_file_count(
                vector_store_id
            )

            remaining_slots = (
                MAX_KNOWLEDGE_FILES - current_count
            )

            if remaining_slots <= 0:

                return (
                    "📚 Your Oddi knowledge base already "
                    "contains 20 files.\n\n"
                    "Please remove an existing file before "
                    "adding another."
                )

            files_to_index = uploaded_files[
                :remaining_slots
            ]

            for uploaded_file in files_to_index:

                try:

                    add_file_to_knowledge_base(
                        uploaded_file,
                        vector_store_id
                    )

                    # Reset the uploaded stream after the knowledge-base
                    # upload so image/document processing can read it again.
                    uploaded_file.stream.seek(0)

                except Exception as e:

                    print(
                        f"Knowledge upload error "
                        f"for {uploaded_file.filename}: {e}"
                    )

                    return (
                        f"⚠️ I couldn't add "
                        f"**{uploaded_file.filename}** "
                        f"to your knowledge base.\n\n"
                        f"Error: {e}"
                    )
    if uploaded_files:
        uploaded_file = uploaded_files[0]
    else:
        uploaded_file = None

    if uploaded_file:
            print("Engine received:", uploaded_file.filename)

    images = []
    documents = ""

    # Audio/video information
    media_transcripts = ""
    video_frames = []
    
    if uploaded_files:

        for uploaded_file in uploaded_files:

            filename = uploaded_file.filename.lower()
            print("Filename:", filename)
            print("Mimetype:", uploaded_file.mimetype)     
            print("Reached ChatGPT section")
            print("Images:", len(images))
            print("Documents length:", len(documents))


            # ---------- IMAGE ----------
            if filename.endswith((".png", ".jpg", ".jpeg", ".webp")):

                image_bytes = uploaded_file.read()

                images.append({
                    "bytes": image_bytes,
                    "mimetype": uploaded_file.mimetype
                })

                print("Image detected")
                print("Image size:", len(image_bytes))
                print("IMAGE BLOCK")
                        # ---------- AUDIO ----------
            elif filename.endswith((
                ".mp3",
                ".wav",
                ".m4a",
                ".aac",
                ".ogg",
                ".flac"
            )):

                print("Audio detected")

                transcript = transcribe_audio(uploaded_file)

                media_transcripts += (
                    "\n\n===== AUDIO: "
                    + uploaded_file.filename
                    + " =====\n"
                )

                media_transcripts += transcript

                print("Audio transcription completed")
                print("AUDIO BLOCK")

                        # ---------- VIDEO ----------
            elif filename.endswith((
                ".mp4",
                ".mov",
                ".avi",
                ".mkv",
                ".webm"
            )):

                print("Video detected")

                with tempfile.NamedTemporaryFile(
                    delete=False,
                    suffix=os.path.splitext(filename)[1]
                ) as temp_video:

                    uploaded_file.save(temp_video.name)
                    video_path = temp_video.name

                try:
                    # ---------- EXTRACT AUDIO ----------
                    audio_path = video_path + ".wav"

                    subprocess.run(
                        [
                            "ffmpeg",
                            "-y",
                            "-i",
                            video_path,
                            "-vn",
                            "-ac",
                            "1",
                            "-ar",
                            "16000",
                            audio_path
                        ],
                        capture_output=True,
                        text=True
                    )

                    if os.path.exists(audio_path):

                        with open(audio_path, "rb") as audio_file:

                            transcript = transcribe_audio(audio_file)

                        media_transcripts += (
                            "\n\n===== VIDEO AUDIO: "
                            + uploaded_file.filename
                            + " =====\n"
                        )

                        media_transcripts += transcript

                    # ---------- EXTRACT VIDEO FRAMES ----------
                    cap = cv2.VideoCapture(video_path)

                    total_frames = int(
                        cap.get(cv2.CAP_PROP_FRAME_COUNT)
                    )

                    fps = cap.get(cv2.CAP_PROP_FPS)

                    if fps <= 0:
                        fps = 25

                    duration = total_frames / fps

                    # Maximum 8 representative frames
                    frame_count = min(8, max(1, int(duration)))

                    for i in range(frame_count):

                        timestamp = (
                            duration * i / frame_count
                        )

                        cap.set(
                            cv2.CAP_PROP_POS_MSEC,
                            timestamp * 1000
                        )

                        success, frame = cap.read()

                        if not success:
                            continue

                        success, encoded = cv2.imencode(
                            ".jpg",
                            frame
                        )

                        if success:

                            video_frames.append({
                                "bytes": encoded.tobytes(),
                                "mimetype": "image/jpeg"
                            })

                    cap.release()

                    print(
                        "Video processing completed:",
                        len(video_frames),
                        "frames"
                    )

                finally:

                    if os.path.exists(video_path):
                        os.remove(video_path)

                    if os.path.exists(audio_path):
                        os.remove(audio_path)

                print("VIDEO BLOCK")
            # ---------- TXT ----------
            elif filename.endswith(".txt"):

                documents += "\n\n===== " + uploaded_file.filename + " =====\n"
                documents += uploaded_file.read().decode("utf-8")

                print("TXT detected")
                print(documents[:200])
                print("TXT BLOCK")

                    # ---------- DOCX ----------
            elif filename.endswith(".docx"):

                print("DOCX detected")

                documents += (
                    "\n\n===== "
                    + uploaded_file.filename
                    + " =====\n"
                )

                docx_text = extract_docx_text(uploaded_file)

                documents += docx_text

                print("DOCX extracted")
                print(documents[:500])
                print("DOCX BLOCK")
            # ---------- PDF ----------
            elif filename.endswith(".pdf"):

                reader = PdfReader(uploaded_file)

                documents += "\n\n===== " + uploaded_file.filename + " =====\n"

                for page in reader.pages:
                    page_text = page.extract_text()

                    if page_text:
                        documents += page_text + "\n"

                print("PDF detected")
                print(documents[:200])
                print("PDF BLOCK")
            # ---------- EXCEL ----------
            elif filename.endswith(".xlsx"):

                workbook = load_workbook(uploaded_file)

                sheet = workbook.active

                documents += "\n\n===== " + uploaded_file.filename + " =====\n"

                for row in sheet.iter_rows(values_only=True):

                    line = " | ".join(str(cell) if cell is not None else "" for cell in row)

                    documents += line + "\n"

                print("EXCEL detected")
                print(documents[:300])
                print("EXCEL BLOCK")

            elif filename.endswith((".pptx", ".ppt")):

                print("PowerPoint detected")

                with tempfile.TemporaryDirectory() as temp_dir:

                    input_path = os.path.join(
                        temp_dir,
                        uploaded_file.filename
                    )

                    uploaded_file.save(input_path)

                    # Convert old .ppt files to .pptx using LibreOffice
                    if filename.endswith(".ppt"):

                        print("Converting .ppt to .pptx using LibreOffice...")

                        soffice_path = shutil.which("soffice")

                        if not soffice_path:
                            soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"

                        if not os.path.exists(soffice_path):
                            raise RuntimeError(
                                "LibreOffice was installed, but soffice.exe could not be found."
                            )

                        result = subprocess.run(
                            [
                                soffice_path,
                                "--headless",
                                "--convert-to",
                                "pptx",
                                "--outdir",
                                temp_dir,
                                input_path
                            ],
                            capture_output=True,
                            text=True
                        )

                        print("LibreOffice output:")
                        print(result.stdout)
                        print(result.stderr)

                        converted_path = os.path.join(
                            temp_dir,
                            os.path.splitext(
                                uploaded_file.filename
                            )[0] + ".pptx"
                        )

                        if not os.path.exists(converted_path):

                            raise RuntimeError(
                                "LibreOffice could not convert the PowerPoint file."
                            )

                        presentation = Presentation(converted_path)

                    else:

                        presentation = Presentation(input_path)

                    documents += (
                        "\n\n===== "
                        + uploaded_file.filename
                        + " =====\n"
                    )

                    for slide_number, slide in enumerate(
                        presentation.slides,
                        start=1
                    ):

                        documents += (
                            f"\n--- Slide {slide_number} ---\n"
                        )

                        for shape in slide.shapes:

                            if hasattr(shape, "text") and shape.text.strip():

                                documents += (
                                    shape.text.strip()
                                    + "\n"
                                )

                    print("POWERPOINT detected")
                    print(documents[:500])
                    print("POWERPOINT BLOCK")
            else:

                print("Unsupported file:", filename)

            # ---------- POWERPOINT ----------

    # 2. Fast Responses
    if not uploaded_files:

        fast_reply = fast_response(user_message)

        if fast_reply:
            return fast_reply


    # Job-context follow-up questions must still work
    # even when job files are uploaded.
    if uploaded_files:

        if is_job_context_question(user_message):

            fast_reply = fast_response(user_message, message)

            if fast_reply:
                return fast_reply
    # 2.5 Resume Intelligence
    resume_analysis = is_resume_analysis_request(user_message)
    if documents:

        content.append({
            "type":"input_text",
            "text":f"""
    Attached documents:

    {documents}
    """
        })
    if media_transcripts:

        content.append({
            "type": "input_text",
            "text": f"""
    Attached audio/video transcription:

    {media_transcripts}

    Use this transcription when answering questions
    about the uploaded audio or video.
    """
        })

    # Add extracted video frames
    for frame in video_frames:

        frame_base64 = base64.b64encode(
            frame["bytes"]
        ).decode("utf-8")

        content.append({
            "type": "input_image",
            "image_url":
            f"data:{frame['mimetype']};base64,{frame_base64}"
        })
    if resume_analysis:

        if not uploaded_files:
            return (
                "📄 **Resume Analysis**\n\n"
                "Please attach your resume first.\n\n"
                "I can analyze **PDF, DOCX, or TXT** resume files."
            )

        if not documents:
            return (
                "⚠️ I received the file, but I couldn't extract readable "
                "resume text from it.\n\n"
                "Please try a PDF, DOCX, or TXT version of your resume."
            )

        user_message = RESUME_ANALYSIS_PROMPT + """

    Here is the extracted resume content:

    """ + documents

    # 1. Commands
    command_reply = handle_command(user_message)
    if command_reply:
        return command_reply


    if user_id is not None:

        if user_id is not None and state.pending_update:
            pending = state.pending_update

            # Make sure this pending update belongs to this user
            if pending.get("user_id") == user_id:

                confirmation = user_message.strip().lower()

                if confirmation in ("y", "yes"):
                    save_memory(
                        user_id,
                        pending["key"],
                        pending["new"]
                    )

                    state.pending_update = None

                    return (
                        f"✅ Updated! I'll remember your "
                        f"{pending['key'].replace('_', ' ')} "
                        f"is '{pending['new']}'."
                    )

                elif confirmation in ("n", "no"):
                    state.pending_update = None

                    return (
                        f"👍 Okay, I'll keep your "
                        f"{pending['key'].replace('_', ' ')} "
                        f"as '{pending['old']}'."
                    )
                memory_reply = recall_memory(user_message, user_id)

                if memory_reply:
                    return memory_reply


    # 4. Save Memory
    if user_id is not None and is_memory_message(user_message):

        return remember(user_message, user_id)

    # 5. Chatgpt
    content = [
        {
            "type": "input_text",
            "text": user_message
        }
    ]
    for image in images:
        print("Building content...")
        image_base64 = base64.b64encode(
            image["bytes"]
        ).decode("utf-8")

        content.append({
            "type": "input_image",
            "image_url":
            f"data:{image['mimetype']};base64,{image_base64}"
        })

    if documents:

        content.append({
            "type":"input_text",
            "text":f"""
    Attached documents:

    {documents}
    """
        })
    print("Creating content...")

    chat_history = []

    if conversation_history:

        for message in conversation_history:

            role = message.get("role")
            text = message.get("text", "")

            if role not in ("user", "assistant"):
                continue

            if not text:
                continue

            chat_history.append({
                "role": role,
                "content": text
            })
    if not chat_history or chat_history[-1]["role"] != "user":
        chat_history.append({
            "role": "user",
            "content": content
        })
    else:
        # Replace the current user's plain text entry
        # with the richer content containing files.
        chat_history[-1] = {
            "role": "user",
            "content": content
        }

    chat_history = chat_history[-100:]

    print("Chat history created")
    print("Messages in memory:", len(chat_history))
    print("Images:", len(images))
    print("Documents:", len(documents))

    print("Sending request to OpenAI...")

    return get_response(
        chat_history,
        vector_store_id
    )